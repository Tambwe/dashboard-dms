from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "generated"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PRIMARY = RGBColor(42, 135, 200)
PRIMARY_LIGHT = RGBColor(232, 244, 251)
SECONDARY = RGBColor(84, 84, 86)
SECONDARY_LIGHT = RGBColor(245, 245, 245)
TERTIARY = RGBColor(157, 72, 56)
ACCENT = RGBColor(212, 140, 116)
DARK = RGBColor(19, 48, 83)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(217, 226, 234)

def add_textbox(slide, left, top, width, height, text, font_size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, fill=None, line_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 8
    tf.margin_right = 8
    tf.margin_top = 6
    tf.margin_bottom = 6
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = color
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line_color is not None:
        box.line.color.rgb = line_color
    return box


def add_box(slide, left, top, width, height, title, bullets, fill, text_color=DARK, accent=None, title_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1.5)

    if accent is not None:
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.12))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.10), width - Inches(0.36), Inches(0.42))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = title_color if fill != SECONDARY_LIGHT else DARK

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.55), width - Inches(0.36), height - Inches(0.7))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.font.size = Pt(11)
        p.font.name = "Aptos"
        p.font.color.rgb = text_color

    return shape


def add_arrow(slide, x1, y1, x2, y2, color=PRIMARY):
    if abs(x2 - x1) >= abs(y2 - y1):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x1, y1 - 10, max(0, x2 - x1), 20)
    else:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x1 - 10, y1, 20, max(0, y2 - y1))
    arrow.line.color.rgb = color
    arrow.line.width = Pt(2.0)
    arrow.fill.background()
    return arrow


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(250, 252, 254)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "DMS Platform – High-Level Architecture"
    run.font.name = "Aptos"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = DARK

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.3))
    stf = subtitle.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Secure access, application services, private connectivity, and mobile data collection"
    sr.font.name = "Aptos"
    sr.font.size = Pt(12)
    sr.font.color.rgb = SECONDARY

    blocks = [
        {
            "title": "1. User Access",
            "fill": PRIMARY_LIGHT,
            "text_color": DARK,
            "accent": PRIMARY,
            "title_color": DARK,
            "pos": (Inches(0.65), Inches(1.3), Inches(2.2), Inches(1.8)),
            "bullets": ["Desktop users", "Admin / analyst roles", "Field and partner access"],
        },
        {
            "title": "2. Identity & SSO",
            "fill": SECONDARY_LIGHT,
            "text_color": DARK,
            "accent": SECONDARY,
            "title_color": DARK,
            "pos": (Inches(3.1), Inches(1.3), Inches(2.3), Inches(1.8)),
            "bullets": ["Azure AD / Okta / Keycloak", "OAuth2 / OIDC", "MFA + RBAC policies"],
        },
        {
            "title": "3. Application Platform",
            "fill": PRIMARY,
            "text_color": WHITE,
            "accent": PRIMARY_LIGHT,
            "title_color": WHITE,
            "pos": (Inches(5.7), Inches(1.15), Inches(2.7), Inches(2.1)),
            "bullets": ["Laravel dashboard", "Reporting and analytics", "Data workflows and APIs"],
        },
        {
            "title": "4. Data & Services",
            "fill": SECONDARY_LIGHT,
            "text_color": DARK,
            "accent": TERTIARY,
            "title_color": DARK,
            "pos": (Inches(8.75), Inches(1.3), Inches(2.2), Inches(1.8)),
            "bullets": ["MySQL / PostgreSQL", "Redis cache and jobs", "GIS, email, alerts"],
        },
        {
            "title": "5. Secure Mobile + Private Network",
            "fill": PRIMARY_LIGHT,
            "text_color": DARK,
            "accent": ACCENT,
            "title_color": DARK,
            "pos": (Inches(2.5), Inches(4.1), Inches(8.5), Inches(1.9)),
            "bullets": ["Offline mobile collection", "HTTPS + TLS 1.2+ encryption", "Private VNet / VPN / firewall / least privilege"],
        },
    ]

    for block in blocks:
        add_box(slide, *block["pos"], block["title"], block["bullets"], block["fill"], block["text_color"], block["accent"], block["title_color"])

    # arrows
    add_arrow(slide, Inches(2.95), Inches(2.2), Inches(3.10), Inches(2.2), PRIMARY)
    add_arrow(slide, Inches(5.45), Inches(2.2), Inches(5.70), Inches(2.2), PRIMARY)
    add_arrow(slide, Inches(8.45), Inches(2.2), Inches(8.75), Inches(2.2), PRIMARY)
    add_arrow(slide, Inches(6.80), Inches(3.25), Inches(6.80), Inches(4.05), ACCENT)

    # footer note
    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.9), Inches(0.4))
    note_tf = note.text_frame
    note_tf.word_wrap = True
    p = note_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Security model: public edge access, authenticated identity, encrypted data transport, and private internal connectivity"
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = SECONDARY

    # small footer brand line
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRIMARY
    strip.line.fill.background()

    prs.save(OUT_DIR / "dms_architecture_boardroom.pptx")
    return prs


def build_pdf():
    pdf_path = OUT_DIR / "dms_architecture_boardroom.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(1600, 900))
    c.setTitle("DMS Platform High-Level Architecture")
    c.setAuthor("Copilot")

    c.setFillColor(HexColor('#F9FBFD'))
    c.rect(0, 0, 1600, 900, stroke=0, fill=1)

    # title
    c.setFillColor(HexColor('#133053'))
    c.setFont('Helvetica-Bold', 30)
    c.drawString(60, 820, 'DMS Platform – High-Level Architecture')

    c.setFillColor(HexColor('#4d5d6b'))
    c.setFont('Helvetica', 14)
    c.drawString(62, 790, 'Secure access, application services, private connectivity, and mobile data collection')

    palette = [
        ('User Access', 60, 560, '#E8F4FB', '#133053', ['Desktop users', 'Admin / analyst roles', 'Field and partner access'], '#2A87C8'),
        ('Identity & SSO', 360, 560, '#F5F5F5', '#133053', ['Azure AD / Okta / Keycloak', 'OAuth2 / OIDC', 'MFA + RBAC'], '#545456'),
        ('Application Platform', 680, 540, '#2A87C8', '#FFFFFF', ['Laravel dashboard', 'Reporting and analytics', 'Data workflows and APIs'], '#2A87C8'),
        ('Data & Services', 1100, 560, '#F5F5F5', '#133053', ['MySQL / PostgreSQL', 'Redis cache and jobs', 'GIS, email, alerts'], '#9d4838'),
        ('Secure Mobile + Private Network', 280, 250, '#E8F4FB', '#133053', ['Offline mobile collection', 'HTTPS + TLS 1.2+', 'Private VNet / VPN / firewall'], '#D48C74'),
    ]

    # draw boxes
    for title, x, y, fill_hex, text_color_hex, lines, accent_hex in palette:
        c.setFillColor(HexColor(fill_hex))
        c.roundRect(x, y, 260, 180, 18, stroke=1, fill=1)
        c.setStrokeColor(HexColor('#D9E2EA'))
        c.setFillColor(HexColor(text_color_hex))
        c.setFont('Helvetica-Bold', 18)
        c.drawString(x + 18, y + 146, title)
        c.setFillColor(HexColor(text_color_hex))
        c.setFont('Helvetica', 11)
        offset = 118
        for line in lines:
            c.drawString(x + 18, y + offset, '• ' + line)
            offset -= 20
        c.setFillColor(HexColor(accent_hex))
        c.rect(x, y + 165, 260, 12, stroke=0, fill=1)

    # connector lines
    c.setStrokeColor(HexColor('#2A87C8'))
    c.setLineWidth(3)
    c.line(320, 650, 360, 650)
    c.line(620, 650, 680, 650)
    c.line(940, 650, 1100, 650)
    c.line(910, 540, 910, 430)

    c.setFillColor(HexColor('#545456'))
    c.setFont('Helvetica-Bold', 12)
    c.drawString(245, 140, 'Security model: public edge access, authenticated identity, encrypted connections, and private internal connectivity')
    c.setFillColor(HexColor('#2A87C8'))
    c.rect(60, 110, 1480, 8, stroke=0, fill=1)
    c.save()


if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print(f"Created: {OUT_DIR / 'dms_architecture_boardroom.pptx'}")
    print(f"Created: {OUT_DIR / 'dms_architecture_boardroom.pdf'}")
